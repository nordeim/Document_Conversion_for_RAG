import gradio as gr
import os
import tempfile
from typing import Tuple, Optional

# For DOCX files
try:
    import docx
    has_docx = True
except ImportError:
    has_docx = False

# For XLSX files
try:
    import openpyxl
    has_openpyxl = True
except ImportError:
    has_openpyxl = False

# For PPTX files
try:
    import pptx
    has_pptx = True
except ImportError:
    has_pptx = False

# For older Office formats
try:
    import win32com.client
    has_win32com = True
except ImportError:
    has_win32com = False

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a .docx file."""
    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from a .xlsx file."""
    wb = openpyxl.load_workbook(file_path, read_only=True)
    full_text = []
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        full_text.append(f"Sheet: {sheet_name}")
        
        for row in sheet.rows:
            row_text = []
            for cell in row:
                if cell.value is not None:
                    row_text.append(str(cell.value))
            if row_text:
                full_text.append('\t'.join(row_text))
    
    return '\n'.join(full_text)

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a .pptx file."""
    prs = pptx.Presentation(file_path)
    full_text = []
    
    for i, slide in enumerate(prs.slides):
        full_text.append(f"Slide {i+1}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    
    return '\n'.join(full_text)

def extract_text_from_doc(file_path: str) -> str:
    """Extract text from a .doc file using win32com."""
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False
    doc = word.Documents.Open(file_path)
    text = doc.Content.Text
    doc.Close()
    word.Quit()
    return text

def extract_text_from_xls(file_path: str) -> str:
    """Extract text from a .xls file using win32com."""
    excel = win32com.client.Dispatch("Excel.Application")
    excel.Visible = False
    wb = excel.Workbooks.Open(file_path)
    full_text = []
    
    for sheet in wb.Sheets:
        full_text.append(f"Sheet: {sheet.Name}")
        used_range = sheet.UsedRange
        for row in range(1, used_range.Rows.Count + 1):
            row_text = []
            for col in range(1, used_range.Columns.Count + 1):
                cell_value = sheet.Cells(row, col).Value
                if cell_value is not None:
                    row_text.append(str(cell_value))
            if row_text:
                full_text.append('\t'.join(row_text))
    
    wb.Close(False)
    excel.Quit()
    return '\n'.join(full_text)

def extract_text_from_ppt(file_path: str) -> str:
    """Extract text from a .ppt file using win32com."""
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(file_path, WithWindow=False)
    full_text = []
    
    for i in range(1, presentation.Slides.Count + 1):
        slide = presentation.Slides.Item(i)
        full_text.append(f"Slide {i}")
        for shape in slide.Shapes:
            if shape.HasTextFrame:
                if shape.TextFrame.HasText:
                    full_text.append(shape.TextFrame.TextRange.Text)
    
    presentation.Close()
    powerpoint.Quit()
    return '\n'.join(full_text)

def convert_office_to_text(file_obj, output_filename: str) -> Tuple[str, Optional[str]]:
    """Convert an Office document to text and save it with the given filename."""
    if not file_obj:
        return "Please upload a file.", None
    
    if not output_filename:
        output_filename = "output.txt"
    
    # Add .txt extension if not present
    if not output_filename.endswith('.txt'):
        output_filename += '.txt'
    
    # Create a temporary file to save the uploaded file
    temp_dir = tempfile.mkdtemp()
    temp_file_path = os.path.join(temp_dir, file_obj.name)
    
    with open(temp_file_path, 'wb') as f:
        f.write(file_obj.read())
    
    # Determine file type and extract text
    file_ext = os.path.splitext(file_obj.name)[1].lower()
    
    try:
        if file_ext == '.docx' and has_docx:
            text = extract_text_from_docx(temp_file_path)
        elif file_ext == '.xlsx' and has_openpyxl:
            text = extract_text_from_xlsx(temp_file_path)
        elif file_ext == '.pptx' and has_pptx:
            text = extract_text_from_pptx(temp_file_path)
        elif file_ext == '.doc' and has_win32com:
            text = extract_text_from_doc(temp_file_path)
        elif file_ext == '.xls' and has_win32com:
            text = extract_text_from_xls(temp_file_path)
        elif file_ext == '.ppt' and has_win32com:
            text = extract_text_from_ppt(temp_file_path)
        else:
            return f"Unsupported file format: {file_ext}", None
    except Exception as e:
        return f"Error processing file: {str(e)}", None
    
    # Save the extracted text to a file
    output_path = os.path.join(temp_dir, output_filename)
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(text)
    
    return text, output_path

# Create Gradio interface
with gr.Blocks(title="Office to Text Converter", theme=gr.themes.Soft()) as app:
    gr.Markdown("# Microsoft Office to Text Converter")
    gr.Markdown("Upload any Microsoft Office document to extract its text content.")
    
    with gr.Row():
        with gr.Column():
            file_input = gr.File(label="Upload Office Document")
            output_filename = gr.Textbox(label="Output Filename", placeholder="output.txt")
            convert_btn = gr.Button("Convert to Text", variant="primary")
        
        with gr.Column():
            text_output = gr.Textbox(label="Extracted Text", lines=15)
            file_output = gr.File(label="Download Text File")
    
    convert_btn.click(
        fn=convert_office_to_text,
        inputs=[file_input, output_filename],
        outputs=[text_output, file_output]
    )
    
    gr.Markdown("""
    ## Supported File Formats
    - Word Documents (.docx, .doc)
    - Excel Spreadsheets (.xlsx, .xls)
    - PowerPoint Presentations (.pptx, .ppt)
    
    ## Required Libraries
    - python-docx (for .docx files)
    - openpyxl (for .xlsx files)
    - python-pptx (for .pptx files)
    - pywin32 (for .doc, .xls, .ppt files - Windows only)
    - gradio
    """)

if __name__ == "__main__":
    app.launch()
  
